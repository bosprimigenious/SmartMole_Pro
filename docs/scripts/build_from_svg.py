#!/usr/bin/env python3
"""SVG → PNG → PPTX pipeline for SmartMole Pro defense deck."""

from __future__ import annotations

import json
import subprocess
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

from svg_rasterize import svg_to_png

DOCS = Path(__file__).resolve().parents[1]
ROOT = DOCS / "ppt-projects" / "smartmole-defense"
SVG_DIR = ROOT / "svg"
PNG_DIR = ROOT / "assets" / "png-from-svg"
OUT = DOCS / "SmartMolePro_答辩PPT.pptx"
OUT_SVG = DOCS / "SmartMolePro_答辩PPT_SVG版.pptx"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def _ensure_svgs():
    if not (SVG_DIR / "manifest.json").is_file():
        print("Generating SVG slides...")
        subprocess.run([sys.executable, str(DOCS / "scripts" / "generate_svg_slides.py")], check=True)


def _convert_all() -> list[Path]:
    manifest = json.loads((SVG_DIR / "manifest.json").read_text(encoding="utf-8"))
    pngs: list[Path] = []
    print("Converting SVG → PNG...")
    for name in manifest:
        svg = SVG_DIR / name
        png = PNG_DIR / name.replace(".svg", ".png")
        if svg_to_png(svg, png):
            print(f"  OK {name} → {png.name}")
            pngs.append(png)
        else:
            raise SystemExit(f"Failed to convert {name}")
    return pngs


def _build_pptx(pngs: list[Path], dest: Path):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]
    for png in pngs:
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(str(png), Inches(0), Inches(0), SLIDE_W, SLIDE_H)
    try:
        prs.save(dest)
        saved = dest
    except PermissionError:
        alt = dest.with_name("_SmartMolePro_SVG_build.pptx")
        prs.save(alt)
        saved = alt
        print(f"WARN: {dest.name} 被占用，已写入 {alt.name}")
    print(f"PPTX: {saved} ({len(pngs)} slides, ~{saved.stat().st_size // 1024} KB)")


def main():
    _ensure_svgs()
    pngs = _convert_all()
    _build_pptx(pngs, OUT_SVG)
    try:
        import shutil

        shutil.copy2(OUT_SVG, OUT)
        print(f"Also copied → {OUT}")
    except OSError as err:
        print(f"Skip copy to {OUT}: {err}")


if __name__ == "__main__":
    main()
