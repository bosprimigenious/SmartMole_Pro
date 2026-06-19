#!/usr/bin/env python3
"""Generate SmartMole Pro defense deck (.pptx) with images & rich layouts.

Design: ppt-master「重庆大学」 #006BB7 / #003D6B / #D4A84B
Run: python scripts/fetch_ppt_assets.py && python scripts/build_defense_ppt.py
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

DOCS = Path(__file__).resolve().parents[1]
OUT = DOCS / "SmartMolePro_答辩PPT.pptx"
ASSETS = DOCS / "ppt-projects" / "smartmole-defense" / "assets" / "images"

C_PRIMARY = RGBColor(0x00, 0x6B, 0xB7)
C_DARK = RGBColor(0x00, 0x3D, 0x6B)
C_GOLD = RGBColor(0xD4, 0xA8, 0x4B)
C_INK = RGBColor(0x1A, 0x2E, 0x44)
C_MUTED = RGBColor(0x6B, 0x7B, 0x8C)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT = RGBColor(0xFA, 0xFC, 0xFF)
C_PANEL = RGBColor(0xF0, 0xF5, 0xFB)
C_OK = RGBColor(0x15, 0x80, 0x3D)
C_WARN = RGBColor(0xC2, 0x41, 0x0C)
C_NO = RGBColor(0xB9, 0x1C, 0x1C)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
FONT = "Microsoft YaHei"


def _img(name: str) -> Path | None:
    p = ASSETS / name
    return p if p.is_file() else None


def _blank(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _rect(slide, l, t, w, h, fill, line=None, line_w=0):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line:
        sh.line.color.rgb = line
        sh.line.width = Pt(line_w)
    else:
        sh.line.fill.background()
    return sh


def _round_rect(slide, l, t, w, h, fill, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line:
        sh.line.color.rgb = line
        sh.line.width = Pt(1.5)
    else:
        sh.line.fill.background()
    return sh


def _textbox(slide, l, t, w, h, text, size=18, bold=False, color=C_INK, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.name = FONT
    p.font.color.rgb = color
    p.alignment = align
    return tb


def _bullets(slide, l, t, w, h, lines, size=16, color=C_INK, line_space=1.15):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(size)
        p.font.name = FONT
        p.font.color.rgb = color
        p.space_after = Pt(5)
        p.line_spacing = line_space
    return tb


def _picture(slide, path: Path, l, t, w, h):
    return slide.shapes.add_picture(str(path), l, t, w, h)


def _framed_image(slide, path: Path | None, l, t, w, h, caption: str | None = None):
    _round_rect(slide, l - Inches(0.06), t - Inches(0.06), w + Inches(0.12), h + Inches(0.12), C_WHITE, C_PRIMARY)
    if path:
        _picture(slide, path, l, t, w, h)
    else:
        _rect(slide, l, t, w, h, C_PANEL, C_MUTED, 1)
        _textbox(slide, l, t + h / 2 - Inches(0.2), w, Inches(0.4), "配图", 14, color=C_MUTED, align=PP_ALIGN.CENTER)
    if caption:
        _textbox(slide, l, t + h + Inches(0.04), w, Inches(0.35), caption, 10, color=C_MUTED, align=PP_ALIGN.CENTER)


def _footer(slide, section: str, num: int):
    _rect(slide, Inches(0), Inches(6.85), SLIDE_W, Inches(0.65), C_LIGHT)
    _textbox(slide, Inches(0.5), Inches(6.95), Inches(4), Inches(0.4), "SmartMole Pro · 第六组", 10, color=C_MUTED)
    _textbox(slide, Inches(4.5), Inches(6.95), Inches(4.3), Inches(0.4), section, 10, color=C_MUTED, align=PP_ALIGN.CENTER)
    badge = _round_rect(slide, Inches(12.2), Inches(6.92), Inches(0.55), Inches(0.38), C_PRIMARY)
    _textbox(slide, Inches(12.2), Inches(6.93), Inches(0.55), Inches(0.35), str(num), 11, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)


def _header_bar(slide, title: str, section: str, num: int):
    _rect(slide, Inches(0), Inches(0), Inches(0.12), SLIDE_H, C_PRIMARY)
    _rect(slide, Inches(0.12), Inches(0), SLIDE_W, Inches(1.05), C_DARK)
    _textbox(slide, Inches(0.55), Inches(0.22), Inches(10), Inches(0.7), title, 28, bold=True, color=C_WHITE)
    _rect(slide, Inches(0.55), Inches(0.95), Inches(1.2), Inches(0.06), C_GOLD)
    _footer(slide, section, num)


def _content_slide(prs, title, bullets, section, num, sub=None, img_name=None, img_w=Inches(5.2)):
    slide = _blank(prs)
    _header_bar(slide, title, section, num)
    has_img = img_name and _img(img_name)
    text_w = Inches(6.3) if has_img else Inches(12.2)
    y = Inches(1.25)
    if sub:
        _textbox(slide, Inches(0.55), y, text_w, Inches(0.45), sub, 14, color=C_MUTED)
        y = Inches(1.65)
    _bullets(slide, Inches(0.55), y, text_w, Inches(5.0), bullets, 16 if has_img else 17)
    if has_img:
        _framed_image(slide, _img(img_name), Inches(7.0), Inches(1.35), img_w, Inches(4.8))
    return slide


def _image_slide(prs, title, img_name, section, num, bullets=None, sub=None, img2=None, gallery=None):
    slide = _blank(prs)
    _header_bar(slide, title, section, num)
    path = _img(img_name)
    if img2 and _img(img2):
        if bullets:
            _bullets(slide, Inches(0.55), Inches(1.2), Inches(12.2), Inches(1.5), bullets, 14)
        _inset(slide, img_name, Inches(0.55), Inches(2.85), Inches(6.0), Inches(3.55), sub)
        _inset(slide, img2, Inches(6.75), Inches(2.85), Inches(6.0), Inches(3.55))
    elif bullets:
        _bullets(slide, Inches(0.55), Inches(1.25), Inches(5.8), Inches(5.2), bullets, 15)
        _framed_image(slide, path, Inches(6.6), Inches(1.2), Inches(6.2), Inches(5.3), sub)
    else:
        _framed_image(slide, path, Inches(0.55), Inches(1.2), Inches(12.2), Inches(5.35), sub)
    if gallery:
        _triple_gallery(slide, gallery, y=Inches(5.05), h=Inches(0.95))
    return slide


def _banner_image(slide, img_name: str, top=Inches(1.08), height=Inches(1.05)):
    path = _img(img_name)
    if path:
        _picture(slide, path, Inches(0.12), top, SLIDE_W - Inches(0.12), height)


def _inset(slide, img_name: str, l, t, w, h, caption=None):
    if _img(img_name):
        _framed_image(slide, _img(img_name), l, t, w, h, caption)


def _triple_gallery(slide, names: list[str], y=Inches(5.55), h=Inches(1.05)):
    w = Inches(4.0)
    gap = Inches(0.25)
    x = Inches(0.55)
    for name in names[:3]:
        if _img(name):
            _framed_image(slide, _img(name), x, y, w, h)
        x += w + gap


def _dual_image_slide(prs, title, img1, img2, section, num, bullets=None, sub=None):
    slide = _blank(prs)
    _header_bar(slide, title, section, num)
    if bullets:
        _bullets(slide, Inches(0.55), Inches(1.2), Inches(12.2), Inches(1.85), bullets, 14)
        top = Inches(3.15)
    else:
        top = Inches(1.2)
    _inset(slide, img1, Inches(0.55), top, Inches(6.0), Inches(3.55), sub)
    _inset(slide, img2, Inches(6.75), top, Inches(6.0), Inches(3.55))
    return slide


def _rich_slide(prs, title, bullets, section, num, main_img, side_img=None, sub=None, gallery=None):
    slide = _blank(prs)
    _header_bar(slide, title, section, num)
    _banner_image(slide, main_img)
    y = Inches(2.25)
    text_w = Inches(6.0) if side_img else Inches(12.2)
    if sub:
        _textbox(slide, Inches(0.55), y, text_w, Inches(0.4), sub, 13, color=C_MUTED)
        y += Inches(0.42)
    _bullets(slide, Inches(0.55), y, text_w, Inches(3.5), bullets, 15)
    if side_img:
        _inset(slide, side_img, Inches(6.85), Inches(2.15), Inches(5.9), Inches(3.55))
    if gallery:
        _triple_gallery(slide, gallery)
    return slide


def _card_grid_slide(prs, title, cards, section, num, hero_img=None, icon_names=None):
    """cards: list of (label, desc, color); icon_names parallel list of png names."""
    slide = _blank(prs)
    _header_bar(slide, title, section, num)
    if hero_img and _img(hero_img):
        _inset(slide, hero_img, Inches(8.15), Inches(1.25), Inches(4.75), Inches(5.35))
    cols, card_w, card_h = 2 if hero_img else 3, Inches(3.75 if hero_img else 3.85), Inches(1.55 if hero_img else 2.15)
    gap_x, gap_y = Inches(0.28), Inches(0.22)
    x0, y0 = Inches(0.55), Inches(1.35)
    for i, (label, desc, color) in enumerate(cards):
        r, c = divmod(i, cols)
        x = x0 + c * (card_w + gap_x)
        y = y0 + r * (card_h + gap_y)
        _round_rect(slide, x, y, card_w, card_h, C_WHITE, C_PRIMARY)
        _rect(slide, x, y, card_w, Inches(0.1), color)
        ix = x + Inches(0.12)
        iy = y + Inches(0.18)
        if icon_names and i < len(icon_names) and _img(icon_names[i]):
            _picture(slide, _img(icon_names[i]), ix, iy, Inches(1.05), Inches(0.65))
            tx = ix + Inches(1.15)
        else:
            tx = ix
        _textbox(slide, tx, y + Inches(0.15), card_w - Inches(0.3), Inches(0.42), label, 14, bold=True, color=C_DARK)
        _textbox(slide, x + Inches(0.15), y + Inches(0.82), card_w - Inches(0.3), card_h - Inches(0.95), desc, 12, color=C_INK)
    return slide


def _table_slide(prs, title, headers, rows, section, num, col_widths=None, img_name=None, img2=None):
    slide = _blank(prs)
    _header_bar(slide, title, section, num)
    table_top = Inches(1.35)
    table_w = Inches(7.8) if img_name and _img(img_name) else Inches(12.4)
    rows_n = len(rows) + 1
    cols_n = len(headers)
    table = slide.shapes.add_table(rows_n, cols_n, Inches(0.45), table_top, table_w, Inches(0.36 * rows_n)).table
    if col_widths:
        scale = 7.8 / 12.4 if img_name and _img(img_name) else 1.0
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w * scale)
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        cell.vertical_anchor = 1
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(11)
            p.font.name = FONT
            p.font.color.rgb = C_WHITE
        cell.fill.solid()
        cell.fill.fore_color.rgb = C_PRIMARY
    for i, row in enumerate(rows, 1):
        bg = C_LIGHT if i % 2 == 0 else C_WHITE
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.text = val
            cell.vertical_anchor = 1
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(10)
                p.font.name = FONT
                p.font.color.rgb = C_INK
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            if j == 2 and val in ("✓", "⚠", "×"):
                color = C_OK if val == "✓" else C_WARN if val == "⚠" else C_NO
                for p in cell.text_frame.paragraphs:
                    p.font.color.rgb = color
                    p.font.bold = True
            if len(headers) == 3 and j == 2 and val.startswith("✓"):
                for p in cell.text_frame.paragraphs:
                    p.font.color.rgb = C_OK
                    p.font.bold = True
    if img_name and _img(img_name):
        _framed_image(slide, _img(img_name), Inches(8.5), Inches(1.35), Inches(4.4), Inches(2.2))
    if img2 and _img(img2):
        _framed_image(slide, _img(img2), Inches(8.5), Inches(3.75), Inches(4.4), Inches(2.2))
    return slide


def _cover_slide(prs):
    slide = _blank(prs)
    cover = _img("cover_arcade.jpg")
    if cover:
        _picture(slide, cover, Inches(0), Inches(0), SLIDE_W, SLIDE_H)
    # 左侧信息栏 + 右侧轻遮罩保证文字可读
    _rect(slide, Inches(0), Inches(0), Inches(4.8), SLIDE_H, C_DARK)
    _rect(slide, Inches(4.75), Inches(0), Inches(8.6), SLIDE_H, C_DARK)  # 半透明感深色
    _rect(slide, Inches(4.7), Inches(0.8), Inches(0.07), Inches(6.0), C_GOLD)
    _textbox(slide, Inches(5.2), Inches(1.8), Inches(7.5), Inches(0.95), "SmartMole Pro", 48, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    _textbox(slide, Inches(5.2), Inches(2.75), Inches(7.5), Inches(0.75), "多模态感知智能打地鼠竞技系统", 26, color=C_GOLD, align=PP_ALIGN.CENTER)
    _textbox(slide, Inches(5.2), Inches(3.55), Inches(7.5), Inches(0.5), "OpenVela (NuttX RTOS) · 课程设计答辩 · 结题版", 16, color=C_LIGHT, align=PP_ALIGN.CENTER)
    _round_rect(slide, Inches(5.5), Inches(4.35), Inches(7.0), Inches(0.55), C_PRIMARY)
    _textbox(slide, Inches(5.5), Inches(4.43), Inches(7.0), Inches(0.4), "综合完成度约 60%  ·  如实标注 ✓ / ⚠ / ×", 14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    _textbox(slide, Inches(0.35), Inches(5.5), Inches(4.2), Inches(1.5),
             "第六组\n张恒基 · 曹佳轩 · 缪钰\n郭志罡 · 张耀辉 · 朱辰骏", 15, color=C_WHITE, align=PP_ALIGN.CENTER)
    return slide


def _closing_slide(prs):
    slide = _blank(prs)
    bg = _img("motivation_arcade.jpg")
    if bg:
        _picture(slide, bg, Inches(0), Inches(0), SLIDE_W, SLIDE_H)
    _rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, C_DARK)
    _textbox(slide, Inches(1), Inches(1.6), Inches(11.3), Inches(0.8), "总结与展望", 40, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    _rect(slide, Inches(5.8), Inches(2.45), Inches(1.7), Inches(0.06), C_GOLD)
    items = [
        ("成果", "五级闯关 + 双板 Wi-Fi 联机可稳定演示，结题文档齐全"),
        ("沉淀", "NuttX 多线程 · UDP 协议 · LVGL UI · wapi WiFi"),
        ("待补", "IP 外置 · 音效全挂钩 · 超声波 · 排行榜 · WS2812B"),
        ("展望", "versus.conf 配置化 · 云端排行 · AI 微调难度"),
    ]
    y = Inches(2.85)
    for label, text in items:
        _round_rect(slide, Inches(1.6), y, Inches(10.1), Inches(0.85), C_WHITE)
        _rect(slide, Inches(1.6), y, Inches(0.08), Inches(0.85), C_GOLD)
        _textbox(slide, Inches(1.85), y + Inches(0.08), Inches(1.2), Inches(0.35), label, 14, bold=True, color=C_PRIMARY)
        _textbox(slide, Inches(3.1), y + Inches(0.1), Inches(8.3), Inches(0.65), text, 15, color=C_INK)
        y += Inches(0.95)
    _textbox(slide, Inches(1), Inches(6.35), Inches(11.3), Inches(0.6), "THANKS · SmartMole Pro 第六组", 24, bold=True, color=C_GOLD, align=PP_ALIGN.CENTER)
    return slide


def build():
    if not ASSETS.is_dir():
        raise SystemExit(f"Assets missing. Run: python scripts/fetch_ppt_assets.py")

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    n = 1

    _cover_slide(prs)
    n += 1

    _rich_slide(prs, "项目背景与动机（起-痛-机-愿）", [
        "起：以课程 WhackMole 基础实验为起点，OpenVela/NuttX 平台可运行",
        "痛：单一触控、无联机竞技、无持久化与声光反馈，产品感不足",
        "机：T113S3 板载 Wi-Fi/音频/GPIO 使多模块集成成为可能",
        "愿：构建可答辩的嵌入式娱乐系统——闯关 + 双板联机为核心",
    ], "背景与目标", n, "hammer_arcade.jpg", "motivation_arcade.jpg",
        gallery=["touch_tablet.jpg", "pcb_board.jpg", "code_dev.jpg"])
    n += 1

    _image_slide(prs, "五层目标体系（结题对照）", "diag_levels.png", "背景与目标", n,
        bullets=["L1 基础 — ✓ 100%", "L2 多模态 — ⚠ 30%", "L3 闯关联机 — ✓ 70%", "L4 声光存储 — ⚠ 50%", "L5 UI WiFi — ⚠ 40%"],
        img2="embedded_board.jpg", gallery=["planning_desk.jpg", "microcontroller.jpg", "sensor_lab.jpg"])
    n += 1

    _card_grid_slide(prs, "已实现技术亮点（结题视角）", [
        ("五级闯关 ✓", "LEVEL 1–5 固定难度参数表", C_OK),
        ("Wi-Fi 联机 ✓", "UDP versus START/SCORE/FINISH", C_OK),
        ("特殊地鼠 ✓", "黄金/炸弹/COMBO 计分逻辑", C_OK),
        ("游戏 UI ✓", "9 洞 + 5 功能按钮", C_OK),
        ("WiFi 菜单 ✓", "wifi_ui.c 图形配置", C_OK),
        ("声光反馈 ⚠", "hit.wav + GPIO LED 降级", C_WARN),
    ], "背景与目标", n, hero_img="mosaic_highlights.jpg",
        icon_names=["icon_level.png", "icon_wifi.png", "icon_mole.png", "icon_ui.png", "icon_menu.png", "icon_av.png"])
    n += 1

    _image_slide(prs, "三层架构（实际实现）", "diag_architecture.png", "系统架构", n,
        bullets=["用户交互层：触摸 + K1 + LVGL", "应用逻辑层：状态机 + storage + versus", "驱动层：NuttX / GPIO / 音频 / wapi"],
        img2="touchscreen_app.jpg", gallery=["code_dev.jpg", "pcb_board.jpg", "embedded_board.jpg"])
    n += 1

    _dual_image_slide(prs, "多线程并发模型", "diag_threads.png", "microcontroller.jpg", "系统架构", n,
        bullets=["LVGL 主线程 · sound_task · led_task · key_task · versus_rx_task", "请求-消费标志位，避免阻塞 GUI"])
    n += 1

    _image_slide(prs, "输入与事件流（当前 vs 规划）", "diag_event_flow.png", "系统架构", n,
        img2="diag_multimodal.png", gallery=["touch_tablet.jpg", "sensor_lab.jpg", "hammer_arcade.jpg"])
    n += 1

    _dual_image_slide(prs, "五级闯关系统 ✓", "diag_mole_ui.png", "diag_level_table.png", "核心模块", n,
        bullets=["LEVEL 按钮循环选关 · 联机共用配置", "黄金 15% / 炸弹 5% · level_configs[] L52–58"])
    n += 1

    _image_slide(prs, "多模态输入 ⚠", "diag_multimodal.png", "核心模块", n,
        bullets=["✓ 触摸 9 洞 + 锤子", "⚠ K1 START", "× K2 / HC-SR04 / 事件队列"],
        img2="touch_tablet.jpg", gallery=["embedded_board.jpg", "sensor_lab.jpg", "hammer_arcade.jpg"])
    n += 1

    _image_slide(prs, "Wi-Fi 双板联机 ✓", "diag_versus.png", "核心模块", n,
        bullets=["versus 模块（张耀辉）", "R528 UDP · START/SCORE/FINISH", "finish_from_peer · 双板联调 · IP 外置", "versus_rx_task 10ms · UI 50ms 刷新", "双板全屏 · 不分屏"],
        img2="wifi_router.jpg", sub="versus_protocol.c · versus_wifi_transport.c",
        gallery=["wifi_network.jpg", "team_meeting.jpg", "code_dev.jpg"])
    n += 1

    _dual_image_slide(prs, "联机 IP 锁定问题与方案", "diag_ip_lock.png", "code_dev.jpg", "核心模块", n,
        bullets=["peer_ip 硬编码 L1103", "换板：改宏 → 编译 → 烧录", "后续：versus.conf / WiFi 弹窗"])
    n += 1

    _dual_image_slide(prs, "声光反馈 ⚠（降级决策）", "diag_sound_led.png", "speaker_audio.jpg", "核心模块", n,
        bullets=["8 wav 仅 hit.wav 挂钩", "GPIO LED 120ms · WS2812B 降级"])
    # 追加 LED 实拍小图
    slide = prs.slides[-1]
    _inset(slide, "led_light.jpg", Inches(11.0), Inches(1.15), Inches(1.95), Inches(1.35))
    n += 1

    _dual_image_slide(prs, "数据存储 ⚠", "diag_storage.png", "database.jpg", "核心模块", n,
        bullets=["whackmole_stats.dat 持久化", "STATS 弹窗 · 未做 littlefs 排行榜"])
    n += 1

    _dual_image_slide(prs, "WiFi 图形连接 ✓", "diag_wifi_ui.png", "wifi_router.jpg", "核心模块", n,
        bullets=["wifi_ui.c LVGL 弹窗", "SCAN / CONNECT · 演示镜像已就绪"])
    n += 1

    _image_slide(prs, "特殊地鼠 + COMBO ✓", "diag_combo.png", "核心模块", n,
        img2="diag_mole_ui.png", gallery=["hammer_arcade.jpg", "cover_arcade.jpg", "motivation_arcade.jpg"])
    n += 1

    _table_slide(prs, "硬件方案与降级矩阵",
        ["模块", "选型", "状态", "降级方案"],
        [
            ["超声波", "HC-SR04", "× 未接", "纯触摸 + K1"],
            ["灯效", "WS2812B", "× 降级", "GPIO LED 闪烁"],
            ["按键", "K1/K2", "⚠ K1", "触摸为主"],
            ["Wi-Fi", "板载 STA", "✓", "wifi_ui 图形配置"],
            ["存储", "文件方案", "⚠", "STATS 弹窗"],
            ["开发板", "DshanPI T113S3", "✓", "—"],
        ], "硬件与成本", n, col_widths=[2.2, 2.5, 1.2, 3.5],
        img_name="microcontroller.jpg", img2="pcb_board.jpg")
    n += 1

    _table_slide(prs, "人员分工（开题方案）",
        ["成员", "负责模块", "备注"],
        [
            ["张恒基", "系统总架构 + 全局集成", "项目负责人"],
            ["曹佳轩", "多模态输入 + 驱动", "外设实现，辅助前端"],
            ["缪钰", "声光 + GUI 重构 + UI", "UI 核心"],
            ["郭志罡", "关卡体系 + AI + 特殊地鼠", "核心框架，主要编码"],
            ["张耀辉", "Wi-Fi 联机 + versus", "versus 核心实现"],
            ["朱辰骏", "存储 + 排行榜 + 成就", "联机联调 + storage API"],
        ], "团队与计划", n, col_widths=[2.0, 4.5, 3.0],
        img_name="team_work.jpg", img2="team_meeting.jpg")
    n += 1

    _dual_image_slide(prs, "协作机制", "diag_git_flow.png", "git_branch.jpg", "团队与计划", n,
        bullets=["versus 核心冻结 · API 接入", "Git feature → review → 合并", "文档四件套 docs/*.typ"])
    n += 1

    _image_slide(prs, "四周里程碑 · 实际达成", "diag_timeline.png", "团队与计划", n,
        img2="planning_desk.jpg", gallery=["team_work.jpg", "code_dev.jpg", "hammer_arcade.jpg"])
    n += 1

    rows = [
        ["音效", "⚠", "hit.wav 已通，7 wav 待挂钩"],
        ["LED", "⚠", "GPIO 降级"],
        ["Wi-Fi 联机", "✓", "versus 已验证"],
        ["五级闯关", "✓", "LEVEL 1–5"],
        ["特殊地鼠/COMBO", "✓", "已实现"],
        ["游戏 UI", "✓", "5 按钮"],
        ["WiFi 菜单", "✓", "wifi_ui.c"],
        ["存储", "⚠", "文件 + STATS"],
        ["HC-SR04", "×", "无代码"],
        ["AI 难度", "×", "固定参数表"],
        ["WS2812B", "×", "已降级"],
        ["排行榜 Tab", "×", "未做"],
    ]
    _table_slide(prs, "开题目标 vs 实际交付（12 项）", ["功能", "状态", "说明"], rows, "完成度对照", n,
        col_widths=[3.0, 0.8, 5.5], img_name="diag_completion.png", img2="database.jpg")
    n += 1

    _image_slide(prs, "技术创新总结 · 与基础实验对比", "diag_innovation.png", "技术创新", n,
        img2="cover_arcade.jpg", gallery=["touch_tablet.jpg", "wifi_router.jpg", "speaker_audio.jpg"])
    n += 1

    _closing_slide(prs)

    try:
        prs.save(OUT)
        saved = OUT
    except PermissionError:
        alt = OUT.with_name("_SmartMolePro_defense_build.pptx")
        prs.save(alt)
        saved = alt
        print(f"WARN: {OUT.name} 被占用，已写入 {alt.name}（请关闭 PowerPoint 后重跑）")
    print(f"Wrote {saved} ({len(prs.slides)} slides, ~{saved.stat().st_size // 1024} KB)")


if __name__ == "__main__":
    build()
